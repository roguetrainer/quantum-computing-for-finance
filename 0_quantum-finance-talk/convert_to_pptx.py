#!/usr/bin/env python3
"""
Convert SLIDES.md to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re

def parse_markdown_slides(md_file):
    """Parse markdown file and extract slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide separators (---)
    slides = content.split('\n---\n')

    parsed_slides = []
    for slide in slides:
        if slide.strip():
            parsed_slides.append(parse_slide_content(slide))

    return parsed_slides

def parse_slide_content(slide_text):
    """Parse individual slide content"""
    lines = slide_text.strip().split('\n')

    slide_data = {
        'title': '',
        'content': [],
        'type': 'content'
    }

    # Extract title (look for ### or ## headers)
    for i, line in enumerate(lines):
        if line.startswith('### Slide'):
            # Title slide format
            match = re.search(r'### Slide \d+: (.+)', line)
            if match:
                slide_data['title'] = match.group(1)
            break
        elif line.startswith('## PART'):
            # Section header
            slide_data['title'] = line.replace('##', '').strip()
            slide_data['type'] = 'section'
            break
        elif line.startswith('### '):
            slide_data['title'] = line.replace('###', '').strip()
            break
        elif line.startswith('# ') and i == 0:
            slide_data['title'] = line.replace('#', '').strip()
            slide_data['type'] = 'title'
            break

    # Extract content (bullet points, text)
    in_visual = False
    in_speaking = False
    content_lines = []

    for line in lines:
        if line.startswith('**Visual'):
            in_visual = True
            continue
        elif line.startswith('**Speaking Notes'):
            in_speaking = True
            continue
        elif line.startswith('**') and (in_visual or in_speaking):
            in_visual = False
            in_speaking = False
            continue

        # Skip metadata and certain sections
        if any(skip in line for skip in ['**Technical Backup', '**Transition:', 'Design Elements:', 'Diagram Specifications:']):
            continue

        # Extract bullet points
        if line.strip().startswith('- ') or line.strip().startswith('• '):
            content_lines.append(line.strip())
        elif line.strip().startswith('│') and '•' in line:
            # Extract from visual layouts
            clean = line.replace('│', '').strip()
            if clean.startswith('•'):
                content_lines.append(clean)

    slide_data['content'] = content_lines
    return slide_data

def create_powerpoint(slides_data, output_file):
    """Create PowerPoint presentation from parsed slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        if not slide_data['title']:
            continue

        if slide_data['type'] == 'title':
            # Title slide
            slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = slide_data['title']
            if slide_data['content']:
                subtitle.text = '\n'.join(slide_data['content'][:3])

        elif slide_data['type'] == 'section':
            # Section header
            slide_layout = prs.slide_layouts[2]  # Section Header layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = slide_data['title']

        else:
            # Content slide
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            title = slide.shapes.title
            title.text = slide_data['title']

            # Add content
            if slide_data['content']:
                if len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    tf.clear()

                    for i, bullet in enumerate(slide_data['content'][:10]):  # Limit to 10 bullets
                        if i == 0:
                            p = tf.paragraphs[0]
                        else:
                            p = tf.add_paragraph()

                        # Clean bullet text
                        text = bullet.replace('- ', '').replace('• ', '')
                        p.text = text
                        p.level = 0
                        p.font.size = Pt(18)

    prs.save(output_file)
    print(f"✓ PowerPoint created: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")

def main():
    md_file = '/Users/ianbuckley/GitHub/quantum-computing-for-finance/0_quantum-finance-talk/SLIDES.md'
    output_file = '/Users/ianbuckley/GitHub/quantum-computing-for-finance/0_quantum-finance-talk/QuantumFinance_Presentation.pptx'

    print("Converting SLIDES.md to PowerPoint...")
    print(f"Input: {md_file}")
    print(f"Output: {output_file}")
    print()

    slides_data = parse_markdown_slides(md_file)
    print(f"Parsed {len(slides_data)} slides from markdown")

    create_powerpoint(slides_data, output_file)
    print()
    print("Done! You can now open the PowerPoint file.")

if __name__ == '__main__':
    main()
